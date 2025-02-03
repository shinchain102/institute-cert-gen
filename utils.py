import os
import pandas as pd
from docx import Document
from pptx import Presentation
from docxtpl import DocxTemplate
import zipfile
import logging
import tempfile
from pathlib import Path
from docx2pdf import convert
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN

def process_template(template_path):
    """Extract variables from template file"""
    try:
        extension = os.path.splitext(template_path)[1].lower()
        text = ""

        if extension == '.docx':
            doc = DocxTemplate(template_path)
            text = '\n'.join([paragraph.text for paragraph in doc.docx.paragraphs])
        elif extension in ['.ppt', '.pptx']:
            prs = Presentation(template_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + '\n'
        else:
            raise ValueError(f"Unsupported file format: {extension}")

        # Extract variables between {{ and }}
        variables = []
        start = 0
        while True:
            start = text.find('{{', start)
            if start == -1:
                break
            end = text.find('}}', start)
            if end == -1:
                break
            variables.append(text[start+2:end].strip())
            start = end + 2
        return list(set(variables))
    except Exception as e:
        logging.error(f"Template processing error: {str(e)}")
        raise

def process_data_file(data_path):
    """Extract column names from data file"""
    try:
        if data_path.endswith('.csv'):
            df = pd.read_csv(data_path)
        else:
            df = pd.read_excel(data_path)
        return df.columns.tolist()
    except Exception as e:
        logging.error(f"Data file processing error: {str(e)}")
        raise

def process_powerpoint_template(template_path, context):
    """Process PowerPoint template with variables"""
    try:
        prs = Presentation(template_path)

        # Process each slide
        for slide in prs.slides:
            # Process each shape in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    # Replace variables in text while preserving formatting
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text = run.text
                            # Replace each variable while preserving the run's formatting
                            for var, value in context.items():
                                placeholder = f"{{{{{var}}}}}"
                                if placeholder in text:
                                    text = text.replace(placeholder, str(value))
                            run.text = text

        # Create temporary file for the processed template
        temp_dir = tempfile.mkdtemp()
        output_path = os.path.join(temp_dir, 'processed_template.pptx')
        prs.save(output_path)
        return output_path
    except Exception as e:
        logging.error(f"PowerPoint template processing error: {str(e)}")
        raise

def convert_to_pdf(input_path, output_dir):
    """Convert document to PDF using appropriate method based on file type"""
    try:
        # Ensure paths are absolute and normalized
        input_path = os.path.abspath(input_path)
        output_dir = os.path.abspath(output_dir)

        # Create output directory if it doesn't exist
        os.makedirs(output_dir, exist_ok=True)

        # Generate output PDF path
        filename = os.path.splitext(os.path.basename(input_path))[0]
        output_path = os.path.join(output_dir, f"{filename}.pdf")

        extension = os.path.splitext(input_path)[1].lower()

        logging.info(f"Converting {input_path} to PDF at {output_path}")

        if extension == '.docx':
            logging.info("Using docx2pdf for Word document conversion")
            try:
                convert(input_path, output_path)
            except Exception as docx_error:
                logging.error(f"docx2pdf conversion error: {str(docx_error)}")
                raise Exception(f"Word to PDF conversion failed: {str(docx_error)}")
        elif extension in ['.ppt', '.pptx']:
            # For PowerPoint files, we'll use LibreOffice for conversion
            try:
                import subprocess
                logging.info("Using LibreOffice for PowerPoint conversion")

                # Construct the LibreOffice command
                cmd = ['soffice', '--headless', '--convert-to', 'pdf', '--outdir', output_dir, input_path]

                # Execute the conversion command
                process = subprocess.Popen(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
                stdout, stderr = process.communicate()

                if process.returncode != 0:
                    error_msg = f"LibreOffice conversion failed: {stderr.decode()}"
                    logging.error(error_msg)
                    raise Exception(error_msg)

                if not os.path.exists(output_path):
                    raise Exception("PDF file was not created by LibreOffice")

            except FileNotFoundError:
                logging.error("LibreOffice not found. Please install LibreOffice.")
                raise Exception("LibreOffice is required for PowerPoint conversion. Please install it on your system.")
        else:
            raise ValueError(f"Unsupported file format: {extension}")

        if not os.path.exists(output_path):
            raise Exception(f"PDF file was not created at {output_path}")

        return output_path

    except Exception as e:
        error_msg = f"Document to PDF conversion error: {str(e)}"
        logging.error(error_msg)
        raise

def generate_certificates(template_path, data_path, mapping):
    """Generate certificates and create zip file"""
    try:
        logging.info(f"Starting certificate generation with template: {template_path}, data: {data_path}")
        logging.info(f"Variable mapping: {mapping}")

        # Read data file
        if data_path.endswith('.csv'):
            df = pd.read_csv(data_path)
        else:
            df = pd.read_excel(data_path)

        # Create temporary directory for generated files
        temp_dir = tempfile.mkdtemp()
        pdf_files = []
        extension = os.path.splitext(template_path)[1].lower()

        for index, row in df.iterrows():
            try:
                # Create context dictionary for template
                context = {var: str(row[col]) for var, col in mapping.items()}

                # Get name and reg-id from the data
                student_name = str(row[mapping['name']]).strip()
                reg_id = str(row[mapping['reg-id']]).strip()

                # Clean filename - replace invalid characters with underscore
                clean_name = "".join(c if c.isalnum() else "_" for c in student_name)
                clean_regid = "".join(c if c.isalnum() else "_" for c in reg_id)

                # Generate filename: name_regid
                base_filename = f"{clean_name}_{clean_regid}"
                temp_output = os.path.join(temp_dir, f"{base_filename}{extension}")

                if extension == '.docx':
                    # Process DOCX template
                    doc = DocxTemplate(template_path)
                    doc.render(context)
                    doc.save(temp_output)
                elif extension in ['.ppt', '.pptx']:
                    # Process PPTX template with improved formatting preservation
                    processed_template = process_powerpoint_template(template_path, context)
                    temp_output = processed_template

                # Convert to PDF
                pdf_output = convert_to_pdf(temp_output, temp_dir)
                # Rename the PDF to match our naming convention
                final_pdf_name = os.path.join(temp_dir, f"{base_filename}.pdf")
                os.rename(pdf_output, final_pdf_name)
                pdf_files.append(final_pdf_name)
                logging.info(f"Generated certificate for {base_filename}")

            except Exception as e:
                logging.error(f"Error generating certificate for row {index}: {str(e)}")
                continue

        if not pdf_files:
            raise ValueError("No certificates were generated successfully")

        # Create zip file
        zip_path = os.path.join(os.path.dirname(template_path), 'certificates.zip')
        with zipfile.ZipFile(zip_path, 'w') as zipf:
            for pdf in pdf_files:
                zipf.write(pdf, os.path.basename(pdf))

        logging.info(f"Created zip file with {len(pdf_files)} certificates")

        # Cleanup temporary files
        for file in os.listdir(temp_dir):
            try:
                os.remove(os.path.join(temp_dir, file))
            except Exception as e:
                logging.warning(f"Error removing temporary file {file}: {str(e)}")

        try:
            os.rmdir(temp_dir)
        except Exception as e:
            logging.warning(f"Error removing temporary directory: {str(e)}")

        return zip_path
    except Exception as e:
        logging.error(f"Certificate generation error: {str(e)}")
        raise